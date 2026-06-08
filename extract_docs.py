import os
import PyPDF2
from pptx import Presentation

docs_dir = r"c:\Users\prita\OneDrive\Documents\Airtel Employee App 2\data docs"
output_file = r"c:\Users\prita\OneDrive\Documents\Airtel Employee App 2\data_docs_extracted.txt"

def extract_pdf(filepath):
    text = ""
    try:
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                # Extract first few pages completely, and then just sample to save space
                page_text = page.extract_text()
                if page_text:
                    if i < 5:
                        text += page_text + "\n"
                    elif i % 2 == 0:
                        text += page_text[:500] + "\n" # Sample parts of later pages
    except Exception as e:
        text += f"Error reading PDF: {e}\n"
    return text

def extract_pptx(filepath):
    text = ""
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text += f"--- Slide {i+1} ---\n" + "\n".join(slide_text) + "\n\n"
    except Exception as e:
        text += f"Error reading PPTX: {e}\n"
    return text

with open(output_file, 'w', encoding='utf-8') as out:
    for filename in os.listdir(docs_dir):
        filepath = os.path.join(docs_dir, filename)
        out.write(f"========== FILE: {filename} ==========\n")
        if filename.endswith('.pdf'):
            out.write(extract_pdf(filepath))
        elif filename.endswith('.pptx'):
            out.write(extract_pptx(filepath))
        out.write("\n\n")

print(f"Extraction complete. Wrote to {output_file}")
